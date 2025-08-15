"""
Servicio mejorado de biblioteca con soporte ampliado para carga de archivos
Versión compacta y optimizada
"""

import os
import io
import json
import hashlib
from typing import Dict, List, Any, Optional
from datetime import datetime
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Verificar bibliotecas disponibles
try:
    import PyPDF2
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    import pandas as pd
    PANDAS_SUPPORT = True
except ImportError:
    PANDAS_SUPPORT = False

try:
    from PIL import Image
    import pytesseract
    OCR_SUPPORT = True
except ImportError:
    OCR_SUPPORT = False


class EnhancedLibraryService:
    """Servicio mejorado de biblioteca con capacidades ampliadas"""
    
    SUPPORTED_EXTENSIONS = {
        # Documentos
        '.txt': 'text/plain',
        '.md': 'text/markdown', 
        '.pdf': 'application/pdf',
        '.doc': 'application/msword',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        # Presentaciones
        '.ppt': 'application/vnd.ms-powerpoint',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        # Hojas de cálculo
        '.xls': 'application/vnd.ms-excel',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.csv': 'text/csv',
        # Imágenes
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.png': 'image/png',
        '.gif': 'image/gif',
        '.bmp': 'image/bmp',
        # Código y datos
        '.json': 'application/json',
        '.xml': 'text/xml',
        '.html': 'text/html',
        '.py': 'text/x-python',
        '.js': 'text/javascript'
    }
    
    SIZE_LIMITS = {
        'default': 50 * 1024 * 1024,  # 50MB
        'image': 10 * 1024 * 1024,    # 10MB
        'document': 100 * 1024 * 1024  # 100MB
    }
    
    def __init__(self):
        self.base_service = None
        try:
            from real_library_service import RealLibraryService
            self.base_service = RealLibraryService()
            logger.info("✅ Servicio base inicializado")
        except:
            logger.warning("⚠️ Servicio base no disponible")
    
    async def upload_document(
        self, 
        file_content: bytes, 
        filename: str, 
        content_type: str = None,
        metadata: Dict[str, Any] = None,
        extract_text: bool = True,
        ocr_enabled: bool = False
    ) -> Dict[str, Any]:
        """Subir y procesar documento con capacidades mejoradas"""
        try:
            # Validar archivo
            file_extension = os.path.splitext(filename)[1].lower()
            if file_extension not in self.SUPPORTED_EXTENSIONS:
                raise ValueError(f"Tipo no soportado: {file_extension}")
            
            # Validar tamaño
            if len(file_content) > self.SIZE_LIMITS['default']:
                raise ValueError(f"Archivo muy grande: {len(file_content) / 1024 / 1024:.2f}MB")
            
            # Extraer contenido si se requiere
            processed_content = ""
            if extract_text:
                processed_content = await self._extract_content(
                    file_content, file_extension, ocr_enabled
                )
            
            # Calcular hash
            file_hash = hashlib.sha256(file_content).hexdigest()
            
            # Metadata completa
            document_metadata = {
                'filename': filename,
                'content_type': content_type or self.SUPPORTED_EXTENSIONS.get(
                    file_extension, 'application/octet-stream'
                ),
                'file_size': len(file_content),
                'file_hash': file_hash,
                'upload_date': datetime.now().isoformat(),
                'text_extracted': bool(processed_content),
                'ocr_used': ocr_enabled and file_extension in ['.jpg', '.jpeg', '.png', '.bmp']
            }
            
            if metadata:
                document_metadata.update(metadata)
            
            # Usar servicio base si está disponible
            if self.base_service:
                try:
                    result = await self.base_service.upload_document(
                        file_content=file_content,
                        filename=filename,
                        content_type=document_metadata['content_type'],
                        metadata=document_metadata
                    )
                    result['enhanced_processing'] = True
                except Exception as e:
                    logger.warning(f"Error usando servicio base: {e}, usando modo independiente")
                    # Fallback a modo independiente
                    document_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file_hash[:8]}"
                    result = {
                        'document_id': document_id,
                        'filename': filename,
                        'processed_content': processed_content,
                        'metadata': document_metadata,
                        'status': 'uploaded',
                        'enhanced_processing': True
                    }
            else:
                # Modo independiente
                document_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file_hash[:8]}"
                result = {
                    'document_id': document_id,
                    'filename': filename,
                    'processed_content': processed_content,
                    'metadata': document_metadata,
                    'status': 'uploaded',
                    'enhanced_processing': True
                }
            
            logger.info(f"✅ Documento procesado: {filename}")
            return result
            
        except Exception as e:
            logger.error(f"❌ Error: {e}")
            raise
    
    async def _extract_content(
        self, 
        file_content: bytes, 
        file_extension: str,
        ocr_enabled: bool = False
    ) -> str:
        """Extraer contenido de texto del archivo"""
        try:
            # Texto plano
            if file_extension in ['.txt', '.md', '.csv']:
                return file_content.decode('utf-8', errors='ignore')
            
            # PDF
            elif file_extension == '.pdf' and PDF_SUPPORT:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
            
            # Word
            elif file_extension in ['.doc', '.docx'] and DOCX_SUPPORT:
                doc = DocxDocument(io.BytesIO(file_content))
                return "\n".join([p.text for p in doc.paragraphs if p.text])
            
            # PowerPoint
            elif file_extension in ['.ppt', '.pptx'] and PPTX_SUPPORT:
                prs = Presentation(io.BytesIO(file_content))
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)
            
            # Excel
            elif file_extension in ['.xls', '.xlsx'] and PANDAS_SUPPORT:
                df = pd.read_excel(io.BytesIO(file_content))
                return df.to_string()
            
            # Imágenes con OCR
            elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp'] and ocr_enabled and OCR_SUPPORT:
                image = Image.open(io.BytesIO(file_content))
                return pytesseract.image_to_string(image, lang='spa+eng')
            
            # JSON
            elif file_extension == '.json':
                data = json.loads(file_content.decode('utf-8'))
                return json.dumps(data, indent=2, ensure_ascii=False)
            
            # Código fuente
            elif file_extension in ['.py', '.js', '.html', '.xml']:
                return file_content.decode('utf-8', errors='ignore')
            
            else:
                return f"[{file_extension}] - Extracción no implementada"
                
        except Exception as e:
            logger.error(f"Error extrayendo contenido: {e}")
            return f"Error: {str(e)}"
    
    async def get_supported_formats(self) -> Dict[str, Any]:
        """Obtener información sobre formatos soportados"""
        return {
            'supported_extensions': list(self.SUPPORTED_EXTENSIONS.keys()),
            'categories': {
                'documents': ['.pdf', '.doc', '.docx', '.txt', '.md'],
                'presentations': ['.ppt', '.pptx'],
                'spreadsheets': ['.xls', '.xlsx', '.csv'],
                'images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp'],
                'code': ['.py', '.js', '.json', '.xml', '.html']
            },
            'features': {
                'ocr_available': OCR_SUPPORT,
                'pdf_support': PDF_SUPPORT,
                'office_support': DOCX_SUPPORT and PPTX_SUPPORT,
                'max_file_size_mb': 100
            },
            'size_limits': {
                k: f"{v / 1024 / 1024:.1f} MB" 
                for k, v in self.SIZE_LIMITS.items()
            }
        }


# Instancia global
enhanced_library_service = EnhancedLibraryService()